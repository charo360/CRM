import asyncio
import base64
import os
import secrets
import uuid
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

PRESENTATIONS_DIR = os.path.join(os.getcwd(), "public", "presentations")
os.makedirs(PRESENTATIONS_DIR, exist_ok=True)

# ── Color Themes ──────────────────────────────────────────────────────────
THEMES = {
    "dark": {
        "bg": RGBColor(0x0B, 0x1D, 0x14),          # deep forest
        "bg_accent": RGBColor(0x07, 0x1A, 0x10),    # darker green
        "accent": RGBColor(0x4C, 0xD1, 0x37),       # brand green
        "accent2": RGBColor(0x00, 0x9B, 0x3A),      # darker brand
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle": RGBColor(0xB0, 0xE0, 0xC0),
        "body": RGBColor(0xD0, 0xE8, 0xD8),
        "light": RGBColor(0xE8, 0xF5, 0xE9),
        "divider": RGBColor(0x4C, 0xD1, 0x37),
    },
    "light": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_accent": RGBColor(0xF0, 0xFD, 0xF4),
        "accent": RGBColor(0x00, 0x9B, 0x3A),
        "accent2": RGBColor(0x00, 0x7A, 0x2E),
        "title": RGBColor(0x0A, 0x26, 0x14),
        "subtitle": RGBColor(0x3A, 0x6B, 0x4A),
        "body": RGBColor(0x33, 0x33, 0x33),
        "light": RGBColor(0xF0, 0xFD, 0xF4),
        "divider": RGBColor(0x00, 0x9B, 0x3A),
    },
    "navy": {
        "bg": RGBColor(0x0F, 0x17, 0x2A),
        "bg_accent": RGBColor(0x14, 0x1E, 0x33),
        "accent": RGBColor(0x38, 0x9C, 0xE6),
        "accent2": RGBColor(0x21, 0x6F, 0xDB),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle": RGBColor(0xA0, 0xC4, 0xE8),
        "body": RGBColor(0xC8, 0xD8, 0xE8),
        "light": RGBColor(0xE0, 0xEE, 0xF8),
        "divider": RGBColor(0x38, 0x9C, 0xE6),
    },
    "warm": {
        "bg": RGBColor(0x1A, 0x12, 0x0B),
        "bg_accent": RGBColor(0x22, 0x18, 0x10),
        "accent": RGBColor(0xF5, 0x9E, 0x0B),
        "accent2": RGBColor(0xD9, 0x77, 0x06),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle": RGBColor(0xFD, 0xD8, 0xA8),
        "body": RGBColor(0xE0, 0xD0, 0xC0),
        "light": RGBColor(0xFF, 0xF7, 0xED),
        "divider": RGBColor(0xF5, 0x9E, 0x0B),
    },
}

DEFAULT_THEME = "dark"

VALID_DECK_STYLES = frozenset({"ribbon", "minimal", "magazine", "split", "spotlight"})
DEFAULT_DECK_STYLE = "ribbon"

# Curated (style, theme) presets — each pair produces a visually distinct deck.
# Different bg colours + different geometry = genuinely unique output every time.
DECK_PRESETS: List[tuple] = [
    ("ribbon",    "dark"),    # deep-green bg, circles, left accent bar
    ("minimal",   "light"),   # white bg, clean centered, thin rule — completely different feel
    ("magazine",  "navy"),    # dark-blue editorial, serif Georgia, right panel
    ("split",     "warm"),    # amber half-hero, earth tones, bold colour block
    ("spotlight", "dark"),    # green bg but content in a big framed card
    ("minimal",   "navy"),    # crisp white-ish navy, sparse geometry
    ("magazine",  "warm"),    # amber editorial serif
    ("split",     "navy"),    # navy half-hero, strong contrast
    ("spotlight", "light"),   # light bg with soft spotlight card
    ("ribbon",    "warm"),    # amber ribbon + circles
]


def normalize_deck_style(name: Optional[str]) -> str:
    if not name or not isinstance(name, str):
        return DEFAULT_DECK_STYLE
    key = name.strip().lower().replace(" ", "_").replace("-", "_")
    return key if key in VALID_DECK_STYLES else DEFAULT_DECK_STYLE


def pick_deck_preset() -> tuple:
    """Return a (deck_style, theme_name) preset pair — both are randomised together
    so successive decks differ in layout AND colour."""
    return secrets.choice(DECK_PRESETS)


def pick_deck_style_random() -> str:
    return secrets.choice(tuple(VALID_DECK_STYLES))


def _fonts_for_deck_style(deck_style: str) -> tuple:
    """(title_font, body_font) readable pairings per layout family."""
    if deck_style == "minimal":
        return "Segoe UI Light", "Segoe UI"
    if deck_style == "magazine":
        return "Georgia", "Calibri"
    if deck_style == "split":
        return "Calibri", "Calibri"
    return "Calibri Light", "Calibri"


def _lighten(color: RGBColor, amount: int = 40) -> RGBColor:
    """Return a lighter version of an RGBColor by adding `amount` to each channel."""
    hex_str = str(color)  # e.g. "009B3A"
    r = min(255, int(hex_str[0:2], 16) + amount)
    g = min(255, int(hex_str[2:4], 16) + amount)
    b = min(255, int(hex_str[4:6], 16) + amount)
    return RGBColor(r, g, b)


# ── Helpers ────────────────────────────────────────────────────────────────

def _add_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text: str, font_size=18,
                 color: RGBColor = RGBColor(0xFF, 0xFF, 0xFF), bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, items: List[str],
                        font_size=16, color: RGBColor = RGBColor(0xFF, 0xFF, 0xFF),
                        bullet_char="●", spacing=Pt(8), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox


# ── Slide Layouts ──────────────────────────────────────────────────────────

def _slide_title(prs, theme, title: str, subtitle: str, deck_style: str = DEFAULT_DECK_STYLE):
    """Title slide — layout varies strongly by deck_style (not only colors)."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, bfont = _fonts_for_deck_style(ds)

    if ds == "minimal":
        # ── MINIMAL title: huge centered title on clean bg, single thin rule ──
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.07), theme["accent"])
        _add_rect(slide, Inches(0), Inches(7.43), Inches(10), Inches(0.07), theme["accent"])
        _add_textbox(slide, Inches(0.7), Inches(2.0), Inches(8.6), Inches(2.0),
                     title, font_size=48, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=tfont)
        _add_textbox(slide, Inches(1.5), Inches(4.15), Inches(7), Inches(0.9),
                     subtitle, font_size=20, color=theme["subtitle"],
                     alignment=PP_ALIGN.CENTER, font_name=bfont)
        _add_rect(slide, Inches(4.0), Inches(5.1), Inches(2.0), Inches(0.07), theme["accent"])
        return

    if ds == "magazine":
        # ── MAGAZINE title: thick left sidebar, oversized serif title right ──
        SIDEBAR = Inches(2.4)
        _add_rect(slide, Inches(0), Inches(0), SIDEBAR, Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.2), Inches(3.2), Inches(2.0), Inches(1.0),
                     "PRESENTED BY", font_size=9, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=bfont)
        _add_textbox(slide, Inches(0.2), Inches(3.85), Inches(2.0), Inches(0.6),
                     subtitle.split("by")[-1].strip() if "by" in subtitle.lower() else "—",
                     font_size=13, color=RGBColor(0xEE, 0xEE, 0xEE),
                     alignment=PP_ALIGN.CENTER, font_name=bfont)
        _add_textbox(slide, SIDEBAR + Inches(0.45), Inches(1.5), Inches(7.0), Inches(3.5),
                     title, font_size=52, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        _add_rect(slide, SIDEBAR + Inches(0.45), Inches(5.25), Inches(6.5), Inches(0.07), theme["divider"])
        _add_textbox(slide, SIDEBAR + Inches(0.45), Inches(5.5), Inches(6.5), Inches(0.7),
                     subtitle, font_size=17, color=theme["body"],
                     alignment=PP_ALIGN.LEFT, font_name=bfont)
        return

    if ds == "split":
        # ── SPLIT title: left 45% = solid accent fill + white title, right = clean ──
        _add_rect(slide, Inches(0), Inches(0), Inches(4.5), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.45), Inches(1.8), Inches(3.7), Inches(3.2),
                     title, font_size=38, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        _add_rect(slide, Inches(0.45), Inches(5.1), Inches(2.8), Inches(0.07),
                  RGBColor(0xFF, 0xFF, 0xFF))
        _add_textbox(slide, Inches(0.45), Inches(5.3), Inches(3.7), Inches(0.9),
                     subtitle, font_size=15, color=RGBColor(0xE8, 0xEE, 0xF2),
                     alignment=PP_ALIGN.LEFT, font_name=bfont)
        _add_rect(slide, Inches(4.55), Inches(0.4), Inches(0.07), Inches(6.7), theme["bg_accent"])
        _add_textbox(slide, Inches(5.0), Inches(3.4), Inches(4.6), Inches(0.7),
                     "DECK", font_size=72, color=theme["bg_accent"], bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        return

    if ds == "spotlight":
        # ── SPOTLIGHT title: full-bleed dark bg, large bright framed card centre ──
        _add_rect(slide, Inches(0.6), Inches(0.9), Inches(8.8), Inches(5.7), theme["bg_accent"])
        _add_rect(slide, Inches(0.6), Inches(0.9), Inches(8.8), Inches(0.15), theme["accent"])
        _add_rect(slide, Inches(0.6), Inches(6.45), Inches(8.8), Inches(0.15), theme["accent"])
        _add_textbox(slide, Inches(1.0), Inches(1.75), Inches(8.0), Inches(2.5),
                     title, font_size=42, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=tfont)
        _add_rect(slide, Inches(3.5), Inches(4.35), Inches(3.0), Inches(0.07), theme["accent"])
        _add_textbox(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.9),
                     subtitle, font_size=20, color=theme["subtitle"],
                     alignment=PP_ALIGN.CENTER, font_name=bfont)
        return

    # ── RIBBON (default): vertical accent bar left, decorative circles ──
    _add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(7.5), theme["accent"])
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(-0.8), Inches(3.5), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _lighten(theme["accent2"], 40)
    shape.line.fill.background()
    shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.6), Inches(5.6), Inches(2.2), Inches(2.2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = _lighten(theme["accent2"], 30)
    shape2.line.fill.background()
    _add_textbox(slide, Inches(0.9), Inches(1.8), Inches(8.0), Inches(1.9),
                 title, font_size=44, color=theme["title"], bold=True, font_name=tfont)
    _add_textbox(slide, Inches(0.9), Inches(3.8), Inches(7.0), Inches(0.9),
                 subtitle, font_size=21, color=theme["subtitle"], font_name=bfont)
    _add_rect(slide, Inches(0.9), Inches(4.85), Inches(3.0), Inches(0.07), theme["accent"])


def _slide_section(prs, theme, title: str, deck_style: str = DEFAULT_DECK_STYLE):
    """Section divider — layout varies by deck_style."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, _ = _fonts_for_deck_style(ds)

    if ds == "minimal":
        _add_textbox(slide, Inches(0.8), Inches(2.85), Inches(8.4), Inches(2.0),
                     title, font_size=40, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=tfont)
        _add_rect(slide, Inches(3.5), Inches(5.0), Inches(3.0), Inches(0.05), theme["accent"])
        return

    if ds == "magazine":
        _add_rect(slide, Inches(0), Inches(4.85), Inches(10), Inches(2.65), theme["accent"])
        _add_textbox(slide, Inches(0.7), Inches(2.1), Inches(8.6), Inches(1.8),
                     title, font_size=38, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        return

    if ds == "split":
        _add_rect(slide, Inches(0), Inches(0), Inches(3.2), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(3.6), Inches(2.9), Inches(5.8), Inches(2.0),
                     title, font_size=34, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        return

    if ds == "spotlight":
        _add_rect(slide, Inches(1.5), Inches(2.2), Inches(7.0), Inches(3.1), theme["bg_accent"])
        _add_rect(slide, Inches(1.5), Inches(2.2), Inches(7.0), Inches(0.08), theme["accent"])
        _add_textbox(slide, Inches(1.8), Inches(2.95), Inches(6.4), Inches(2.0),
                     title, font_size=32, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=tfont)
        return

    # ribbon + default
    _add_rect(slide, Inches(1.2), Inches(2.0), Inches(7.6), Inches(3.5), theme["bg_accent"])
    _add_rect(slide, Inches(1.2), Inches(2.0), Inches(7.6), Inches(0.08), theme["accent"])
    _add_textbox(slide, Inches(1.5), Inches(2.8), Inches(7), Inches(2.0),
                 title, font_size=36, color=theme["title"], bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=tfont)


def _slide_content(prs, theme, title: str, content: List[str], deck_style: str = DEFAULT_DECK_STYLE):
    """Content slide — header + bullets; geometry changes dramatically by deck_style."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, bfont = _fonts_for_deck_style(ds)

    if ds == "minimal":
        # ── MINIMAL: white/light feel, big bold title left-flush, clean rule, no chrome ──
        # Large title block at top with lots of breathing room
        _add_textbox(slide, Inches(0.65), Inches(0.45), Inches(9.0), Inches(1.1),
                     title, font_size=32, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.65), Inches(1.45), Inches(9.0), Inches(0.055), theme["accent"])
        _add_bullet_textbox(slide, Inches(0.85), Inches(1.75), Inches(8.65), Inches(5.2),
                            content, font_size=18, color=theme["body"], bullet_char="·",
                            spacing=Pt(14), font_name=bfont)

    elif ds == "magazine":
        # ── MAGAZINE: left sidebar = thick colour column, content right — editorial ──
        SIDEBAR = Inches(2.4)
        _add_rect(slide, Inches(0), Inches(0), SIDEBAR, Inches(7.5), theme["accent"])
        # Title rotated 90° not possible in pptx easily — place it in sidebar area
        _add_textbox(slide, Inches(0.18), Inches(1.0), Inches(2.1), Inches(5.5),
                     title, font_size=22, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        # Main content right of sidebar
        _add_bullet_textbox(slide, SIDEBAR + Inches(0.35), Inches(0.45), Inches(7.0), Inches(6.6),
                            content, font_size=18, color=theme["body"], bullet_char="—",
                            spacing=Pt(12), font_name=bfont)
        # Thin top accent rule over content area
        _add_rect(slide, SIDEBAR + Inches(0.35), Inches(0.42), Inches(7.0), Inches(0.05), theme["divider"])

    elif ds == "split":
        # ── SPLIT: top 40% = full-width accent band with title; bottom 60% = bullets ──
        BAND = Inches(2.8)
        _add_rect(slide, Inches(0), Inches(0), Inches(10), BAND, theme["accent"])
        _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), BAND - Inches(0.8),
                     title, font_size=34, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        _add_bullet_textbox(slide, Inches(0.7), BAND + Inches(0.25), Inches(8.6), Inches(4.2),
                            content, font_size=18, color=theme["body"], bullet_char="▸",
                            spacing=Pt(11), font_name=bfont)

    elif ds == "spotlight":
        # ── SPOTLIGHT: full-width framed card, title inside card header ──
        _add_rect(slide, Inches(0.3), Inches(0.25), Inches(9.4), Inches(6.95), theme["bg_accent"])
        _add_rect(slide, Inches(0.3), Inches(0.25), Inches(9.4), Inches(1.25), theme["accent"])
        _add_textbox(slide, Inches(0.65), Inches(0.4), Inches(8.8), Inches(1.0),
                     title, font_size=28, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        _add_bullet_textbox(slide, Inches(0.65), Inches(1.65), Inches(8.8), Inches(5.15),
                            content, font_size=18, color=theme["body"], bullet_char="▶",
                            spacing=Pt(11), font_name=bfont)

    else:
        # ── RIBBON (default): thin top accent bar, left-flush title + underline ──
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), theme["accent"])
        _add_textbox(slide, Inches(0.7), Inches(0.3), Inches(8.5), Inches(0.9),
                     title, font_size=30, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.22), Inches(2.2), Inches(0.06), theme["accent"])
        _add_bullet_textbox(slide, Inches(0.7), Inches(1.55), Inches(8.5), Inches(5.3),
                            content, font_size=18, color=theme["body"], bullet_char="▸",
                            spacing=Pt(10), font_name=bfont)


def _slide_two_column(prs, theme, title: str, left_items: List[str], right_items: List[str],
                      left_header: str = "", right_header: str = "", deck_style: str = DEFAULT_DECK_STYLE):
    """Two-column content slide."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, bfont = _fonts_for_deck_style(ds)
    use_split_header = ds == "split"
    use_minimal_header = ds == "minimal"
    use_mag = ds == "magazine"
    use_spot = ds == "spotlight"

    if use_split_header:
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.95), theme["accent"])
        _add_textbox(slide, Inches(0.55), Inches(0.18), Inches(8.5), Inches(0.65),
                     title, font_size=26, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, font_name=tfont)
        title_y = Inches(1.25)
    elif use_minimal_header:
        _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=26, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.12), Inches(8.0), Inches(0.04), theme["accent"])
        title_y = Inches(1.35)
    elif use_mag:
        _add_rect(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.45), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=28, color=theme["title"], bold=True, font_name=tfont)
        title_y = Inches(1.3)
    elif use_spot:
        _add_rect(slide, Inches(0.45), Inches(0.35), Inches(9.1), Inches(6.45), theme["bg_accent"])
        _add_rect(slide, Inches(0.45), Inches(0.35), Inches(9.1), Inches(0.08), theme["accent"])
        _add_textbox(slide, Inches(0.75), Inches(0.52), Inches(8.5), Inches(0.75),
                     title, font_size=27, color=theme["title"], bold=True, font_name=tfont)
        title_y = Inches(1.25)
    else:
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), theme["accent"])
        _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=28, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.2), Inches(1.8), Inches(0.05), theme["accent"])
        title_y = Inches(1.6)

    col_top = title_y + Inches(0.15)
    _add_rect(slide, Inches(0.5), col_top, Inches(4.3), Inches(5.0), theme["bg_accent"])
    _add_rect(slide, Inches(5.2), col_top, Inches(4.3), Inches(5.0), theme["bg_accent"])

    y = col_top + Inches(0.2)
    if left_header:
        _add_textbox(slide, Inches(0.7), y, Inches(3.9), Inches(0.5),
                     left_header, font_size=16, color=theme["accent"], bold=True)
        y = y + Inches(0.55)

    _add_bullet_textbox(slide, Inches(0.7), y, Inches(3.9), Inches(4.0),
                        left_items, font_size=15, color=theme["body"], bullet_char="▸",
                        spacing=Pt(8), font_name=bfont)

    y = col_top + Inches(0.2)
    if right_header:
        _add_textbox(slide, Inches(5.4), y, Inches(3.9), Inches(0.5),
                     right_header, font_size=16, color=theme["accent"], bold=True)
        y = y + Inches(0.55)

    _add_bullet_textbox(slide, Inches(5.4), y, Inches(3.9), Inches(4.0),
                        right_items, font_size=15, color=theme["body"], bullet_char="▸",
                        spacing=Pt(8), font_name=bfont)

    _add_rect(slide, Inches(4.95), col_top + Inches(0.15), Inches(0.03), Inches(4.5), theme["divider"])


def _slide_quote(prs, theme, quote: str, attribution: str = "", deck_style: str = DEFAULT_DECK_STYLE):
    """Quote / highlight slide — large text; framing depends on deck_style."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    qfont = "Georgia" if ds in ("magazine", "ribbon", "spotlight") else _fonts_for_deck_style(ds)[1]
    bfont = _fonts_for_deck_style(ds)[1]

    if ds == "minimal":
        _add_textbox(slide, Inches(1.0), Inches(2.0), Inches(8), Inches(2.8),
                     quote, font_size=24, color=theme["title"], bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=qfont)
        if attribution:
            _add_textbox(slide, Inches(1.0), Inches(4.85), Inches(8), Inches(0.6),
                         f"— {attribution}", font_size=15, color=theme["subtitle"],
                         alignment=PP_ALIGN.CENTER, font_name=bfont)
        return

    if ds == "split":
        _add_rect(slide, Inches(0), Inches(5.1), Inches(10), Inches(2.4), theme["accent"])
        _add_textbox(slide, Inches(0.9), Inches(1.35), Inches(8.2), Inches(3.2),
                     quote, font_size=24, color=theme["title"], bold=False,
                     alignment=PP_ALIGN.LEFT, font_name=qfont)
        if attribution:
            _add_textbox(slide, Inches(0.9), Inches(4.55), Inches(8), Inches(0.5),
                         f"— {attribution}", font_size=15, color=theme["subtitle"],
                         alignment=PP_ALIGN.LEFT, font_name=bfont)
        return

    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(2), Inches(2),
                 "❝", font_size=72, color=theme["accent"], bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Georgia")
    _add_textbox(slide, Inches(1.2), Inches(2.2), Inches(7.5), Inches(2.5),
                 quote, font_size=26, color=theme["title"], bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Georgia")
    if attribution:
        _add_textbox(slide, Inches(1.2), Inches(4.8), Inches(7.5), Inches(0.6),
                     f"— {attribution}", font_size=16, color=theme["subtitle"],
                     alignment=PP_ALIGN.LEFT, font_name=bfont)
    _add_rect(slide, Inches(1.2), Inches(5.6), Inches(3), Inches(0.05), theme["accent"])


def _slide_image_text(prs, theme, title: str, content: List[str], image_url: Optional[str] = None,
                      deck_style: str = DEFAULT_DECK_STYLE):
    """Content slide with image placeholder on right, text on left."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, bfont = _fonts_for_deck_style(ds)
    body_top = Inches(1.55)

    if ds == "split":
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.95), theme["accent"])
        _add_textbox(slide, Inches(0.55), Inches(0.18), Inches(8.5), Inches(0.65),
                     title, font_size=26, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, font_name=tfont)
        body_top = Inches(1.2)
    elif ds == "minimal":
        _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=26, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.12), Inches(7.5), Inches(0.04), theme["accent"])
        body_top = Inches(1.38)
    elif ds == "magazine":
        _add_rect(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.45), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=28, color=theme["title"], bold=True, font_name=tfont)
        body_top = Inches(1.32)
    elif ds == "spotlight":
        _add_rect(slide, Inches(0.45), Inches(0.35), Inches(9.1), Inches(6.45), theme["bg_accent"])
        _add_rect(slide, Inches(0.45), Inches(0.35), Inches(9.1), Inches(0.08), theme["accent"])
        _add_textbox(slide, Inches(0.75), Inches(0.52), Inches(8.5), Inches(0.75),
                     title, font_size=27, color=theme["title"], bold=True, font_name=tfont)
        body_top = Inches(1.22)
    else:
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), theme["accent"])
        _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=28, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.2), Inches(1.8), Inches(0.05), theme["accent"])
        body_top = Inches(1.6)

    _add_bullet_textbox(slide, Inches(0.7), body_top, Inches(4.8), Inches(5.0),
                        content, font_size=17, color=theme["body"], bullet_char="▸",
                        spacing=Pt(10), font_name=bfont)

    img_top = body_top
    _add_rect(slide, Inches(5.8), img_top, Inches(3.8), Inches(4.5), theme["bg_accent"],
              line_color=theme["divider"])
    _add_textbox(slide, Inches(5.8), img_top + Inches(1.75), Inches(3.8), Inches(0.8),
                 "📷  Image", font_size=16, color=theme["subtitle"],
                 alignment=PP_ALIGN.CENTER)

    if image_url:
        try:
            slide.shapes.add_picture(image_url, Inches(5.9), img_top + Inches(0.05), Inches(3.6), Inches(4.3))
        except Exception:
            pass


def _slide_key_points(prs, theme, title: str, points: List[Dict[str, str]],
                      deck_style: str = DEFAULT_DECK_STYLE):
    """Key points / icon grid slide — up to 4 points with title + description each."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, bfont = _fonts_for_deck_style(ds)
    grid_y = Inches(1.6)

    if ds == "split":
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.95), theme["accent"])
        _add_textbox(slide, Inches(0.55), Inches(0.18), Inches(8.5), Inches(0.65),
                     title, font_size=26, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, font_name=tfont)
        grid_y = Inches(1.2)
    elif ds == "minimal":
        _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=26, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.12), Inches(8.0), Inches(0.04), theme["accent"])
        grid_y = Inches(1.38)
    elif ds == "magazine":
        _add_rect(slide, Inches(0), Inches(0), Inches(0.2), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.45), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=28, color=theme["title"], bold=True, font_name=tfont)
        grid_y = Inches(1.32)
    elif ds == "spotlight":
        _add_rect(slide, Inches(0.45), Inches(0.35), Inches(9.1), Inches(6.45), theme["bg_accent"])
        _add_rect(slide, Inches(0.45), Inches(0.35), Inches(9.1), Inches(0.08), theme["accent"])
        _add_textbox(slide, Inches(0.75), Inches(0.52), Inches(8.5), Inches(0.75),
                     title, font_size=27, color=theme["title"], bold=True, font_name=tfont)
        grid_y = Inches(1.22)
    else:
        _add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), theme["accent"])
        _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8.5), Inches(0.8),
                     title, font_size=28, color=theme["title"], bold=True, font_name=tfont)
        _add_rect(slide, Inches(0.7), Inches(1.2), Inches(1.8), Inches(0.05), theme["accent"])

    positions = [
        (Inches(0.5), grid_y),
        (Inches(5.2), grid_y),
        (Inches(0.5), grid_y + Inches(2.6)),
        (Inches(5.2), grid_y + Inches(2.6)),
    ]
    card_w = Inches(4.3)
    card_h = Inches(2.2)

    for i, point in enumerate(points[:4]):
        px, py = positions[i]
        # Card background
        _add_rect(slide, px, py, card_w, card_h, theme["bg_accent"])
        # Card left accent
        _add_rect(slide, px, py, Inches(0.08), card_h, theme["accent"])
        # Point title
        pt_title = point.get("title", point.get("header", ""))
        pt_desc = point.get("description", point.get("body", point.get("content", "")))
        if pt_title:
            _add_textbox(slide, px + Inches(0.3), py + Inches(0.2), card_w - Inches(0.5), Inches(0.5),
                         pt_title, font_size=18, color=theme["accent"], bold=True)
        if pt_desc:
            _add_textbox(slide, px + Inches(0.3), py + Inches(0.8), card_w - Inches(0.5), Inches(1.2),
                         pt_desc, font_size=14, color=theme["body"], font_name=bfont)


def _slide_ending(prs, theme, business_name: str, tagline: str = "", deck_style: str = DEFAULT_DECK_STYLE):
    """Thank-you / closing slide — matches deck_style family."""
    ds = normalize_deck_style(deck_style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, theme["bg"])
    tfont, bfont = _fonts_for_deck_style(ds)

    if ds == "minimal":
        _add_textbox(slide, Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.2),
                     "Thank You", font_size=42, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=tfont)
        _add_textbox(slide, Inches(0.8), Inches(3.85), Inches(8.4), Inches(0.7),
                     business_name, font_size=20, color=theme["subtitle"],
                     alignment=PP_ALIGN.CENTER, font_name=bfont)
        if tagline:
            _add_textbox(slide, Inches(0.8), Inches(4.65), Inches(8.4), Inches(0.6),
                         tagline, font_size=15, color=theme["body"],
                         alignment=PP_ALIGN.CENTER, font_name=bfont)
        _add_rect(slide, Inches(3.5), Inches(5.35), Inches(3.0), Inches(0.05), theme["accent"])
        return

    if ds == "magazine":
        _add_rect(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.55), Inches(2.35), Inches(8.5), Inches(1.1),
                     "Thank You", font_size=40, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        _add_textbox(slide, Inches(0.55), Inches(3.65), Inches(8), Inches(0.7),
                     business_name, font_size=21, color=theme["subtitle"], font_name=bfont)
        if tagline:
            _add_textbox(slide, Inches(0.55), Inches(4.45), Inches(8), Inches(0.6),
                         tagline, font_size=15, color=theme["body"], font_name=bfont)
        return

    if ds == "split":
        _add_rect(slide, Inches(0), Inches(0), Inches(4.25), Inches(7.5), theme["accent"])
        _add_textbox(slide, Inches(0.45), Inches(2.55), Inches(3.6), Inches(1.1),
                     "Thank You", font_size=34, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=tfont)
        _add_textbox(slide, Inches(4.75), Inches(2.35), Inches(4.9), Inches(0.85),
                     business_name, font_size=24, color=theme["title"], bold=False, font_name=bfont)
        if tagline:
            _add_textbox(slide, Inches(4.75), Inches(3.35), Inches(4.9), Inches(0.65),
                         tagline, font_size=16, color=theme["subtitle"], font_name=bfont)
        _add_rect(slide, Inches(4.55), Inches(5.5), Inches(5.2), Inches(0.08), theme["bg_accent"])
        return

    if ds == "spotlight":
        _add_rect(slide, Inches(1.1), Inches(1.85), Inches(7.8), Inches(3.8), theme["bg_accent"])
        _add_rect(slide, Inches(1.1), Inches(1.85), Inches(7.8), Inches(0.1), theme["accent"])
        _add_textbox(slide, Inches(1.4), Inches(2.25), Inches(7.2), Inches(1.0),
                     "Thank You", font_size=36, color=theme["title"], bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=tfont)
        _add_textbox(slide, Inches(1.4), Inches(3.45), Inches(7.2), Inches(0.65),
                     business_name, font_size=20, color=theme["subtitle"],
                     alignment=PP_ALIGN.CENTER, font_name=bfont)
        if tagline:
            _add_textbox(slide, Inches(1.4), Inches(4.2), Inches(7.2), Inches(0.55),
                         tagline, font_size=15, color=theme["body"],
                         alignment=PP_ALIGN.CENTER, font_name=bfont)
        return

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-1), Inches(4), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _lighten(theme["accent2"], 30)
    shape.line.fill.background()
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), theme["accent"])
    _add_textbox(slide, Inches(0.8), Inches(2.2), Inches(8), Inches(1.2),
                 "Thank You", font_size=44, color=theme["title"], bold=True,
                 alignment=PP_ALIGN.LEFT, font_name=tfont)
    _add_rect(slide, Inches(0.8), Inches(3.5), Inches(2.5), Inches(0.06), theme["accent"])
    _add_textbox(slide, Inches(0.8), Inches(3.9), Inches(7), Inches(0.7),
                 business_name, font_size=22, color=theme["subtitle"], font_name=bfont)
    if tagline:
        _add_textbox(slide, Inches(0.8), Inches(4.7), Inches(7), Inches(0.6),
                     tagline, font_size=16, color=theme["body"], font_name=bfont)


# ── Thumbnail Generator ──────────────────────────────────────────────────────

def _rgb_tuple(color: RGBColor) -> tuple:
    """Convert RGBColor to (r, g, b) tuple for Pillow."""
    h = str(color)  # e.g. "0B1D14"
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _thumb_text_center_x(draw, text: str, font, W: int) -> int:
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    return max(0, (W - tw) // 2)


def _generate_thumbnail(
    prs,
    theme,
    title: str,
    business_name: str,
    thumb_path: str,
    deck_style: str = DEFAULT_DECK_STYLE,
):
    """Generate a PNG thumbnail approximating the title slide for the chosen deck_style."""
    ds = normalize_deck_style(deck_style)
    W, H = 1280, 720
    bg_accent = _rgb_tuple(theme["bg_accent"])
    bg = _rgb_tuple(theme["bg"])
    accent = _rgb_tuple(theme["accent"])
    accent2 = _rgb_tuple(theme["accent2"])
    title_col = _rgb_tuple(theme["title"])
    subtitle_col = _rgb_tuple(theme["subtitle"])

    img = Image.new("RGB", (W, H), bg)
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("calibril.ttf", 56)
    except Exception:
        font_title = ImageFont.truetype("arial.ttf", 56) if os.name == "nt" else ImageFont.load_default()
    try:
        font_sub = ImageFont.truetype("calibri.ttf", 28)
    except Exception:
        font_sub = ImageFont.truetype("arial.ttf", 28) if os.name == "nt" else ImageFont.load_default()

    sub = f"Presented by {business_name}"

    def _wrap(text: str, max_chars: int = 36) -> str:
        if len(text) <= max_chars:
            return text
        words, line, lines = text.split(), "", []
        for w in words:
            if len(line) + len(w) + 1 > max_chars:
                lines.append(line.strip())
                line = w + " "
            else:
                line += w + " "
        lines.append(line.strip())
        return "\n".join(lines[:3])

    if ds == "minimal":
        # Clean: top + bottom rule, huge centered title
        draw.rectangle([0, 0, W, 9], fill=accent)
        draw.rectangle([0, H - 9, W, H], fill=accent)
        tx = _thumb_text_center_x(draw, title[:36], font_title, W)
        draw.text((max(40, tx), 195), _wrap(title, 32), fill=title_col, font=font_title)
        sx = _thumb_text_center_x(draw, sub, font_sub, W)
        draw.text((max(40, sx), 390), sub, fill=subtitle_col, font=font_sub)
        # Small accent rule under title
        draw.rectangle([int(W * 0.38), 460, int(W * 0.62), 467], fill=accent)

    elif ds == "magazine":
        # Thick left sidebar
        SIDE = int(W * 0.24)
        draw.rectangle([0, 0, SIDE, H], fill=accent)
        draw.text((50, 200), _wrap(title, 20), fill=title_col, font=font_title)
        draw.text((SIDE + 40, 310), sub, fill=subtitle_col, font=font_sub)
        draw.rectangle([SIDE + 40, 570, SIDE + 40 + 400, 577], fill=accent)

    elif ds == "split":
        # Left 45% solid accent, right = bg
        split_w = int(W * 0.45)
        draw.rectangle([0, 0, split_w, H], fill=accent)
        draw.text((50, 195), _wrap(title, 18), fill=(255, 255, 255), font=font_title)
        draw.rectangle([50, 510, 340, 519], fill=(255, 255, 255))
        draw.text((50, 535), sub[:55], fill=(232, 238, 242), font=font_sub)
        draw.rectangle([split_w + 28, 55, split_w + 35, 665], fill=bg_accent)

    elif ds == "spotlight":
        # Big framed card on background
        m = 55
        draw.rectangle([m, m, W - m, H - m], fill=bg_accent)
        draw.rectangle([m, m, W - m, m + 16], fill=accent)
        draw.rectangle([m, H - m - 16, W - m, H - m], fill=accent)
        tx = _thumb_text_center_x(draw, title[:36], font_title, W)
        draw.text((max(m + 30, tx), 175), _wrap(title, 32), fill=title_col, font=font_title)
        draw.rectangle([int(W * 0.35), 440, int(W * 0.65), 449], fill=accent)
        sx = _thumb_text_center_x(draw, sub, font_sub, W)
        draw.text((max(m + 30, sx), 465), sub, fill=subtitle_col, font=font_sub)

    else:
        # Ribbon: left bar + circles
        draw.rectangle([0, 0, 18, H], fill=accent)
        lighter = tuple(min(255, c + 40) for c in accent2)
        draw.ellipse([W - 320, -90, W + 230, 430], fill=lighter)
        lighter2 = tuple(min(255, c + 30) for c in accent2)
        draw.ellipse([-70, H - 210, 160, H + 50], fill=lighter2)
        draw.text((110, 220), _wrap(title, 32), fill=title_col, font=font_title)
        draw.text((110, 400), sub, fill=subtitle_col, font=font_sub)
        draw.rectangle([110, H - 175, 420, H - 167], fill=accent)

    # Slide count badge bottom-right
    slide_count = len(prs.slides)
    badge_text = f"{slide_count} slides"
    try:
        font_badge = ImageFont.truetype("calibri.ttf", 18)
    except Exception:
        font_badge = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), badge_text, font=font_badge)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    bx, by = W - tw - 40, H - th - 20
    draw.rounded_rectangle([bx - 12, by - 8, bx + tw + 12, by + th + 8], radius=8, fill=accent)
    draw.text((bx, by), badge_text, fill=(255, 255, 255), font=font_badge)

    img.save(thumb_path, "PNG")


# ── Layout Router ──────────────────────────────────────────────────────────

LAYOUT_MAP = {
    "title": _slide_title,
    "section": _slide_section,
    "content": _slide_content,
    "two_column": _slide_two_column,
    "quote": _slide_quote,
    "image_text": _slide_image_text,
    "key_points": _slide_key_points,
    "ending": _slide_ending,
}


def generate_presentation(
    title: str,
    slides_data: List[Dict[str, Any]],
    business_name: str = "My Business",
    theme_name: str = DEFAULT_THEME,
    tagline: str = "",
    deck_style: str = DEFAULT_DECK_STYLE,
) -> Dict[str, Any]:
    """
    Generates a professional PowerPoint presentation (.pptx).

    deck_style: visual layout family (orthogonal to theme colors): ribbon (default bar+circles),
      minimal (clean centered), magazine (editorial split + Georgia), split (half accent hero),
      spotlight (framed card). Invalid values fall back to ribbon.

    slides_data format — each slide object:
      Common keys:
        - layout (str): "title" | "section" | "content" | "two_column" |
                        "quote" | "image_text" | "key_points" | "ending"
                        Default: "content" (first slide auto-gets "title")

      Layout-specific keys:
        title:   title, subtitle
        section: title
        content: title, content (list of str)
        two_column: title, left_items, right_items, left_header?, right_header?
        quote:   quote, attribution?
        image_text: title, content (list of str), image_url?
        key_points: title, points (list of {title, description})
        ending:  (uses business_name + tagline — no extra keys needed)
    """
    try:
        theme = THEMES.get(theme_name, THEMES[DEFAULT_THEME])
        ds = normalize_deck_style(deck_style)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for idx, sd in enumerate(slides_data):
            layout = sd.get("layout", "content").lower()

            # Auto-make first slide a title slide if not specified
            if idx == 0 and layout not in ("title", "section", "ending"):
                layout = "title"

            # Auto-add ending slide if last slide isn't already one
            if idx == len(slides_data) - 1 and layout not in ("ending", "title"):
                # First render the current slide
                _render_slide(prs, theme, layout, sd, business_name, tagline, ds)
                # Then append ending
                _slide_ending(prs, theme, business_name, tagline, ds)
                continue

            _render_slide(prs, theme, layout, sd, business_name, tagline, ds)

        # If no slides at all, add title + ending
        if not slides_data:
            _slide_title(prs, theme, title, f"Presented by {business_name}", ds)
            _slide_ending(prs, theme, business_name, tagline, ds)

        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(PRESENTATIONS_DIR, filename)
        prs.save(filepath)

        download_url = f"/api/media/presentations/{filename}"

        # Generate PNG thumbnail preview of the title slide
        thumb_filename = f"preview_{uuid.uuid4().hex[:8]}.png"
        thumb_path = os.path.join(PRESENTATIONS_DIR, thumb_filename)
        thumb_url = f"/api/media/presentations/{thumb_filename}"
        _generate_thumbnail(prs, theme, title, business_name, thumb_path, deck_style=ds)

        return {
            "success": True,
            "filename": filename,
            "filepath": filepath,
            "url": download_url,
            "thumbnail_url": thumb_url,
            "deck_style": ds,
        }

    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"error": str(e)}


async def generate_presentation_with_upload(
    title: str,
    slides_data: List[Dict[str, Any]],
    business_name: str = "My Business",
    style: Optional[Dict[str, Any]] = None,
    theme_name: str = DEFAULT_THEME,
    tagline: str = "",
    deck_style: str = DEFAULT_DECK_STYLE,
) -> Dict[str, Any]:
    """Generate presentation and upload to S3. Returns S3 URL."""
    result = generate_presentation(
        title,
        slides_data,
        business_name,
        theme_name=theme_name,
        tagline=tagline,
        deck_style=deck_style,
    )
    if result.get("error"):
        return result

    filepath = result["filepath"]
    file_url = None

    try:
        from pathlib import Path
        from image_handler import S3Handler

        file_bytes = Path(filepath).read_bytes()
        b64 = base64.b64encode(file_bytes).decode()
        s3_name = f"pptx-{uuid.uuid4().hex[:8]}.pptx"
        file_url = await S3Handler.upload_file(
            b64, s3_name, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        import logging
        logging.warning("[presentation_service] S3 upload failed: %s", e)
    finally:
        try:
            os.unlink(filepath)
        except Exception:
            pass

    return {
        "success": True,
        "url": file_url or result.get("url"),
        "thumbnail_url": result.get("thumbnail_url"),
        "filename": result.get("filename"),
        "deck_style": result.get("deck_style"),
    }


# ── Visual (Gemini-image) Presentation ────────────────────────────────────

async def _fetch_image_bytes(url: str) -> Optional[bytes]:
    """Download image from a public URL. Returns None on failure."""
    try:
        import httpx
        async with httpx.AsyncClient(timeout=60.0) as client:
            r = await client.get(url)
            r.raise_for_status()
            return r.content
    except Exception as e:
        import logging
        logging.warning("[visual_pptx] image fetch failed %s: %s", url, e)
        return None


_SLIDE_IMAGE_QUALITY_SUFFIX = (
    " — ultra-clean professional background for a business presentation slide, widescreen 16:9. "
    "Real-world cinematic photography: soft depth of field, natural light, calm and breathable. "
    "Muted mid-tones so white text reads clearly on top. "
    "NO geometric shapes, NO waves, NO abstract patterns, NO gradients, NO UI elements, "
    "NO text, NO watermarks, NO logos, NO borders anywhere in the image."
)

_DESIGNED_SLIDE_QUALITY_SUFFIX = (
    " — Render this as a complete, flat presentation slide image, widescreen 16:9 ratio. "
    "Style: premium Canva/Google Slides quality — clean sans-serif fonts (Inter, Montserrat, or Helvetica Neue), "
    "bold title at top, bullet points below with generous line spacing. "
    "Every single word must be PERFECTLY SPELLED exactly as given — render each letter precisely, no hallucinated words. "
    "Ample whitespace, refined color palette, tasteful brand accents (subtle, not loud). "
    "No laptop frames, no device mockups, no browser windows, no screen chrome, no UI borders — just the flat slide. "
    "DO NOT add decorative AI-template elements: no random blobs, no memphis shapes, no neon/glow, no gradient mesh, "
    "no dense pattern overlays, no fake icon packs, no stickers, no emojis, no 3D renders. "
    "DO NOT invent extra text. Use ONLY the provided TITLE/BODY strings. "
    "The design must look like it came from a human designer at a professional design agency (bespoke), not a template tool."
)

_SLIDE_STYLE_MAP = {
    "title":   "Sweeping cinematic landscape — mountain range, modern city skyline, or vast ocean at golden hour. "
               "Epic scale, dramatic natural lighting, deep depth of field, open negative space. "
               "NO dollar signs, NO currency symbols, NO icons, NO floating symbols, NO charts, NO UI elements, NO data visualizations.",
    "content": "Clean architectural interior or calm natural scene — marble surface, glass office, forest path. "
               "Soft natural light, muted tones, breathable — nothing busy or artificial.",
    "data":    "Dark polished surface — black marble, dark wood desk, or night city reflection. "
               "Minimal, authoritative, premium feel. No geometric shapes, no patterns.",
    "closing": "Wide open panoramic — sunrise over city, coastline at dusk, or aerial landscape. "
               "Warm amber and deep blue natural tones, inspiring and aspirational.",
}


async def _gen_slide_image(
    slide_prompt: str,
    brand_color: str = "",
    logo_url: Optional[str] = None,
    quality: str = "pro",
    slide_role: str = "content",
    ai_designed: bool = True,
) -> Optional[str]:
    """Generate one landscape Gemini image for a slide at award-winning quality."""
    try:
        if ai_designed:
            full_prompt = f"{slide_prompt}{_DESIGNED_SLIDE_QUALITY_SUFFIX}"
        else:
            style_hint = _SLIDE_STYLE_MAP.get(slide_role, _SLIDE_STYLE_MAP["content"])
            full_prompt = f"{style_hint} {slide_prompt}{_SLIDE_IMAGE_QUALITY_SUFFIX}"
        from nano_banana_service import generate_creative_image
        result = await generate_creative_image(
            prompt=full_prompt,
            format="landscape",
            quality=quality,
            brand_color=brand_color,
            logo_url=logo_url,
            brand_color_mode="subtle" if ai_designed else "dominant",
        )
        return result.get("image_url") if result.get("success") else None
    except Exception as e:
        import logging
        logging.warning("[visual_pptx] slide image gen failed: %s", e)
        return None


def _add_visual_slide(
    prs: "Presentation",
    image_bytes: Optional[bytes],
    title: str,
    body_lines: List[str],
    brand_color_hex: str = "",
    is_title_slide: bool = False,
    ai_designed: bool = False,
) -> None:
    """Add one slide: full-bleed background image + text overlay."""
    import io as _io
    from pptx.util import Inches as _In, Pt as _Pt, Emu as _Emu
    from pptx.dml.color import RGBColor as _RGB
    from pptx.enum.text import PP_ALIGN as _ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    W, H = prs.slide_width, prs.slide_height

    # ── Background ──
    if image_bytes:
        try:
            pic = slide.shapes.add_picture(_io.BytesIO(image_bytes), 0, 0, W, H)
            # Send image to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        except Exception:
            pass  # fall through to solid bg if image fails

    if ai_designed and image_bytes:
        # Gemini rendered the complete slide — no pptx elements needed
        return

    # Semi-transparent dark overlay using a shape (pptx doesn't support true opacity
    # on fills, so we use a very dark near-black with alpha via theme trick — use
    # a solid rectangle at 55% lightness equivalent via dark colour)
    overlay_color = _RGB(0x0A, 0x0A, 0x0A)
    ov = slide.shapes.add_shape(1, 0, 0, W, H)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    ov.fill.solid()
    ov.fill.fore_color.rgb = overlay_color
    ov.line.fill.background()
    # Set transparency via XML (python-pptx doesn't expose alpha directly)
    try:
        from lxml import etree as _et
        solidFill = ov.fill._xPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                alpha = _et.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha.set('val', '50000')  # 50% opaque → 50% transparent — lets background breathe
    except Exception:
        pass

    # ── Brand accent bar (top) ──
    try:
        r = int(brand_color_hex[1:3], 16) if brand_color_hex and len(brand_color_hex) == 7 else 0x4C
        g = int(brand_color_hex[3:5], 16) if brand_color_hex and len(brand_color_hex) == 7 else 0xD1
        b = int(brand_color_hex[5:7], 16) if brand_color_hex and len(brand_color_hex) == 7 else 0x37
        accent = _RGB(r, g, b)
    except Exception:
        accent = _RGB(0x4C, 0xD1, 0x37)

    bar = slide.shapes.add_shape(1, 0, 0, W, _In(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # ── Title text ──
    title_top = _In(1.6) if not is_title_slide else _In(2.2)
    title_size = 40 if is_title_slide else 32
    title_box = slide.shapes.add_textbox(_In(0.7), title_top, _In(8.6), _In(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = _Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = _RGB(0xFF, 0xFF, 0xFF)
    p.font.name = "Calibri Light"
    p.alignment = _ALIGN.LEFT

    # ── Accent rule under title ──
    rule_top = title_top + _In(1.55)
    rule = slide.shapes.add_shape(1, _In(0.7), rule_top, _In(2.2), _In(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = accent
    rule.line.fill.background()

    # ── Body bullets ──
    if body_lines:
        body_top = rule_top + _In(0.22)
        body_box = slide.shapes.add_textbox(_In(0.7), body_top, _In(8.6), _In(2.8))
        btf = body_box.text_frame
        btf.word_wrap = True
        for i, line in enumerate(body_lines[:6]):
            bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            bp.text = f"  •  {line}"
            bp.font.size = _Pt(18)
            bp.font.color.rgb = _RGB(0xE8, 0xE8, 0xE8)
            bp.font.name = "Calibri"


async def create_visual_presentation_async(
    topic: str,
    slides_plan: List[Dict[str, Any]],
    business_name: str = "My Business",
    brand_color: str = "",
    logo_url: Optional[str] = None,
    quality: str = "fast",
    ai_designed: bool = True,
) -> Dict[str, Any]:
    """
    Build a PPTX where every slide has a **Gemini-generated full-bleed background image**
    with the slide title and bullet points overlaid on a semi-transparent dark panel (or
    fully AI-designed slide images if ai_designed=True).

    slides_plan — list of dicts, each with:
        title (str)        : slide headline
        body  (list[str])  : 2-5 bullet points / key lines
        image_prompt (str) : what to generate for the slide background image
        is_title (bool)    : True for the cover slide (default False except slide 0)
    """
    import logging as _log

    # Step 1 — generate all slide images in parallel (concurrency capped at 4)
    sem = asyncio.Semaphore(4)

    async def _bounded_gen(sd: Dict[str, Any]) -> Optional[str]:
        async with sem:
            # Determine slide role for style hint/prompt building
            if sd.get("is_title"):
                role = "title"
            elif any(kw in (sd.get("title") or "").lower() for kw in ("thank", "next step", "contact", "cta", "call to action", "conclusion")):
                role = "closing"
            elif any(kw in (sd.get("title") or "").lower() for kw in ("data", "metric", "result", "number", "stat", "revenue", "growth", "kpi")):
                role = "data"
            else:
                role = "content"

            raw = (sd.get("image_prompt") or sd.get("image_concept") or "").strip()
            if ai_designed:
                title = (sd.get("title") or "").strip()
                body_lines = sd.get("body") or []
                bullets_text = ""
                if body_lines:
                    bullets_text = "Bullets: " + ", ".join(f"'{b}'" for b in body_lines)

                bullets_str = "\n".join(f"• {b}" for b in body_lines) if body_lines else ""

                # Keep logo usage minimal to avoid distortion artifacts.
                use_logo = logo_url if role in ("title", "closing") else None

                if role == "title":
                    prompt = (
                        f"Design a complete widescreen 16:9 presentation COVER slide for '{business_name}'.\n"
                        f"TITLE TEXT (render exactly, perfectly spelled): \"{title}\"\n"
                        f"SUBTITLE TEXT (render exactly): \"{business_name}\"\n"
                        f"Background: {raw or 'dramatic cinematic cityscape or landscape at golden hour — real photography only, no effects'}.\n"
                        f"Layout: left-aligned grid. Large bold title. Subtitle below. "
                        f"Keep 15–20% clear margin around text; do not crowd edges.\n"
                        f"Typography: ONE font family only; 2 weights max (bold + regular). "
                        f"STRICTLY NO: glowing lines, network connections, light trails, digital overlays, tech grid effects, "
                        f"particle effects, neon glow, animated shine, geometric decorations, watermarks, device mockups, "
                        f"UI chrome, floating icons, charts, emoji, dollar signs."
                    )
                elif role == "closing":
                    prompt = (
                        f"Design a complete widescreen 16:9 presentation CLOSING slide.\n"
                        f"TITLE TEXT (render exactly, perfectly spelled): \"{title}\"\n"
                        f"BODY TEXT (render exactly):\n{bullets_str}\n"
                        f"Background: {raw or 'wide panoramic sunrise or calm coastline at dusk — real photography only, no effects'}.\n"
                        f"Layout: left-aligned. Bold title, body text below with generous line spacing. "
                        f"Open, breathable, high contrast.\n"
                        f"Typography: ONE font family only; 2 weights max. "
                        f"STRICTLY NO: glowing lines, network connections, light trails, digital overlays, tech grid effects, "
                        f"particle effects, neon glow, geometric decorations, watermarks, device mockups, UI chrome."
                    )
                elif role == "data":
                    prompt = (
                        f"Design a complete widescreen 16:9 presentation DATA/METRICS slide.\n"
                        f"TITLE TEXT (render exactly, perfectly spelled): \"{title}\"\n"
                        f"BODY TEXT (render exactly):\n{bullets_str}\n"
                        f"Background: {raw or 'dark polished surface — black marble or dark wood. Real photography only, no effects, no patterns'}.\n"
                        f"Layout: bold title at top, stats/bullets below with clear hierarchy and generous spacing. "
                        f"Minimal, authoritative, numbers easy to read.\n"
                        f"Typography: ONE font family only; 2 weights max. "
                        f"STRICTLY NO: glowing lines, network connections, light trails, digital overlays, tech grid effects, "
                        f"particle effects, neon glow, geometric decorations, watermarks, device mockups, UI chrome."
                    )
                else:  # content
                    prompt = (
                        f"Design a complete widescreen 16:9 presentation CONTENT slide.\n"
                        f"TITLE TEXT (render exactly, perfectly spelled): \"{title}\"\n"
                        f"BODY TEXT (render exactly, each bullet on its own line):\n{bullets_str}\n"
                        f"Background: {raw or 'clean architectural interior or calm natural scene — real photography only, no effects'}.\n"
                        f"Layout: left-aligned. Bold title at top, bullet points below with generous line spacing. "
                        f"Ample whitespace, high contrast text, never cluttered. Max 3 bullet points.\n"
                        f"Typography: ONE font family only; 2 weights max. "
                        f"STRICTLY NO: glowing lines, network connections, light trails, digital overlays, tech grid effects, "
                        f"particle effects, neon glow, animated shine, geometric decorations, watermarks, device mockups, UI chrome."
                    )
                return await _gen_slide_image(prompt, brand_color=brand_color, logo_url=use_logo, quality=quality, slide_role=role, ai_designed=True)
            else:
                if not raw:
                    raw = (
                        f"Cinematic real-world background for a professional business presentation slide "
                        f"titled '{sd.get('title', topic)}'. "
                        f"Natural photography — architecture, landscape, or calm environment. "
                        f"No geometric shapes, no waves, no abstract patterns, minimal and sophisticated."
                    )
                return await _gen_slide_image(raw, brand_color=brand_color, logo_url=None, quality=quality, slide_role=role, ai_designed=False)

    image_urls: List[Optional[str]] = await asyncio.gather(
        *[_bounded_gen(sd) for sd in slides_plan]
    )

    # Step 2 — download image bytes for all slides in parallel
    async def _dl(url: Optional[str]) -> Optional[bytes]:
        return await _fetch_image_bytes(url) if url else None

    image_bytes_list: List[Optional[bytes]] = await asyncio.gather(
        *[_dl(u) for u in image_urls]
    )

    # Step 3 — assemble PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # widescreen 16:9
    prs.slide_height = Inches(7.5)

    for i, (sd, img_bytes) in enumerate(zip(slides_plan, image_bytes_list)):
        is_title = sd.get("is_title", i == 0)
        _add_visual_slide(
            prs,
            image_bytes=img_bytes,
            title=sd.get("title", f"Slide {i + 1}"),
            body_lines=sd.get("body") or [],
            brand_color_hex=brand_color,
            is_title_slide=bool(is_title),
            ai_designed=ai_designed,
        )

    # Step 4 — save locally then upload to S3
    filename = f"visual-pptx-{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(PRESENTATIONS_DIR, filename)
    prs.save(filepath)

    file_url = f"/api/media/presentations/{filename}"
    try:
        from pathlib import Path
        from image_handler import S3Handler
        file_bytes = Path(filepath).read_bytes()
        b64 = base64.b64encode(file_bytes).decode()
        s3_name = f"visual-pptx-{uuid.uuid4().hex[:8]}.pptx"
        s3_url = await S3Handler.upload_file(
            b64, s3_name,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        if s3_url:
            file_url = s3_url
    except Exception as e:
        _log.warning("[visual_pptx] S3 upload skipped: %s", e)
    finally:
        try:
            os.unlink(filepath)
        except Exception:
            pass

    slide_count = len(slides_plan)
    generated_count = sum(1 for u in image_urls if u)
    return {
        "success": True,
        "url": file_url,
        "filename": filename,
        "slide_count": slide_count,
        "images_generated": generated_count,
        "topic": topic,
        "ai_designed": ai_designed,
        # Return slides + image URLs so the caller can pass them back for per-slide regeneration
        "slides": slides_plan,
        "image_urls": [u or "" for u in image_urls],
    }


async def regenerate_single_slide_async(
    slides_plan: List[Dict[str, Any]],
    image_urls: List[str],
    slide_index: int,
    instruction: str,
    brand_color: str = "",
    quality: str = "pro",
    topic: str = "",
    logo_url: Optional[str] = None,
    ai_designed: bool = False,
) -> Dict[str, Any]:
    """
    Regenerate ONE slide's background image based on new instructions,
    then rebuild and re-upload the full .pptx with that image swapped in.

    Returns the same shape as create_visual_presentation_async.
    """
    import logging as _log

    if slide_index < 0 or slide_index >= len(slides_plan):
        return {"error": f"slide_index {slide_index} is out of range (deck has {len(slides_plan)} slides)."}

    sd = slides_plan[slide_index]

    # Build a refined prompt: original image_prompt + user instruction
    original_prompt = (sd.get("image_prompt") or sd.get("image_concept") or "").strip()
    if instruction.strip():
        if original_prompt:
            refined_prompt = f"{original_prompt}. Modification: {instruction.strip()}"
        else:
            refined_prompt = instruction.strip()
    else:
        refined_prompt = original_prompt or f"Abstract professional background for: {sd.get('title', 'slide')}"

    # Determine slide role
    if sd.get("is_title") or slide_index == 0:
        role = "title"
    elif any(kw in (sd.get("title") or "").lower() for kw in ("thank", "next step", "contact", "cta", "conclusion")):
        role = "closing"
    elif any(kw in (sd.get("title") or "").lower() for kw in ("data", "metric", "result", "stat", "revenue", "growth", "kpi")):
        role = "data"
    else:
        role = "content"

    if ai_designed:
        title = (sd.get("title") or "").strip()
        body_lines = sd.get("body") or []
        bullets_text = ""
        if body_lines:
            bullets_text = "Bullets: " + ", ".join(f"'{b}'" for b in body_lines)

        if role == "title":
            prompt = (
                f"A complete, professionally designed presentation COVER slide. "
                f"Title: '{title}'. Subtitle/business: '{topic}'. "
                f"Style/Background/Changes: {refined_prompt}. "
                f"Elegant corporate layout, bold title, readable high-contrast typography, brand accents."
            )
        elif role == "closing":
            prompt = (
                f"A complete, professionally designed presentation CLOSING/CTA slide. "
                f"Title: '{title}'. Details: '{', '.join(body_lines)}'. "
                f"Style/Background/Changes: {refined_prompt}. "
                f"Strong CTA layout, clean typography, readable content, brand elements."
            )
        elif role == "data":
            prompt = (
                f"A complete, professionally designed presentation DATA/METRIC slide. "
                f"Title: '{title}'. Data: '{', '.join(body_lines)}'. "
                f"Style/Background/Changes: {refined_prompt}. "
                f"Premium slide highlighting key stats or figures, clean fonts, balanced negative space."
            )
        else: # content
            prompt = (
                f"A complete, professionally designed presentation CONTENT slide. "
                f"Title: '{title}'. {bullets_text}. "
                f"Style/Background/Changes: {refined_prompt}. "
                f"Widescreen layout with bold title at top, clean readable bullet points, brand accents."
            )
        new_url = await _gen_slide_image(prompt, brand_color=brand_color, logo_url=logo_url, quality=quality, slide_role=role, ai_designed=True)
    else:
        new_url = await _gen_slide_image(refined_prompt, brand_color=brand_color, logo_url=None, quality=quality, slide_role=role, ai_designed=False)

    if not new_url:
        return {"error": "Image generation failed for this slide. Please try again."}

    # Build updated image_urls list with the new URL swapped in
    updated_urls: List[Optional[str]] = [u if u else None for u in image_urls]
    updated_urls[slide_index] = new_url

    # Download all image bytes in parallel
    async def _dl(url: Optional[str]) -> Optional[bytes]:
        return await _fetch_image_bytes(url) if url else None

    image_bytes_list: List[Optional[bytes]] = await asyncio.gather(*[_dl(u) for u in updated_urls])

    # Also update the slide's image_prompt so it reflects the new version
    updated_slides = [dict(s) for s in slides_plan]
    updated_slides[slide_index] = {**updated_slides[slide_index], "image_prompt": refined_prompt}

    # Rebuild full .pptx
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, (sd_item, img_bytes) in enumerate(zip(updated_slides, image_bytes_list)):
        is_title = sd_item.get("is_title", i == 0)
        _add_visual_slide(
            prs,
            image_bytes=img_bytes,
            title=sd_item.get("title", f"Slide {i + 1}"),
            body_lines=sd_item.get("body") or [],
            brand_color_hex=brand_color,
            is_title_slide=bool(is_title),
            ai_designed=ai_designed,
        )

    filename = f"visual-pptx-{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(PRESENTATIONS_DIR, filename)
    prs.save(filepath)

    file_url = f"/api/media/presentations/{filename}"
    try:
        from pathlib import Path
        from image_handler import S3Handler
        file_bytes = Path(filepath).read_bytes()
        b64 = base64.b64encode(file_bytes).decode()
        s3_name = f"visual-pptx-{uuid.uuid4().hex[:8]}.pptx"
        s3_url = await S3Handler.upload_file(
            b64, s3_name,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        if s3_url:
            file_url = s3_url
    except Exception as e:
        _log.warning("[visual_pptx] S3 re-upload skipped: %s", e)
    finally:
        try:
            os.unlink(filepath)
        except Exception:
            pass

    return {
        "success": True,
        "url": file_url,
        "filename": filename,
        "slide_count": len(updated_slides),
        "images_generated": sum(1 for u in updated_urls if u),
        "regenerated_slide_index": slide_index,
        "regenerated_slide_title": sd.get("title", f"Slide {slide_index + 1}"),
        "topic": topic,
        "ai_designed": ai_designed,
        "slides": updated_slides,
        "image_urls": [u or "" for u in updated_urls],
    }


def _render_slide(
    prs,
    theme,
    layout: str,
    sd: Dict[str, Any],
    business_name: str,
    tagline: str,
    deck_style: str,
):
    ds = normalize_deck_style(deck_style)
    if layout == "title":
        _slide_title(
            prs,
            theme,
            sd.get("title", "Presentation"),
            sd.get("subtitle", f"Presented by {business_name}"),
            ds,
        )
    elif layout == "section":
        _slide_section(prs, theme, sd.get("title", "Section"), ds)
    elif layout == "two_column":
        _slide_two_column(
            prs,
            theme,
            sd.get("title", ""),
            sd.get("left_items", sd.get("left", [])),
            sd.get("right_items", sd.get("right", [])),
            sd.get("left_header", ""),
            sd.get("right_header", ""),
            ds,
        )
    elif layout == "quote":
        _slide_quote(
            prs,
            theme,
            sd.get("quote", sd.get("content", "")),
            sd.get("attribution", ""),
            ds,
        )
    elif layout == "image_text":
        _slide_image_text(
            prs,
            theme,
            sd.get("title", ""),
            sd.get("content", []),
            sd.get("image_url"),
            ds,
        )
    elif layout == "key_points":
        _slide_key_points(prs, theme, sd.get("title", ""), sd.get("points", []), ds)
    elif layout == "ending":
        _slide_ending(prs, theme, business_name, tagline, ds)
    else:
        _slide_content(prs, theme, sd.get("title", "Slide"), sd.get("content", []), ds)


# ── Clean Deck Design System ────────────────────────────────────────────────────
# Entirely rendered via python-pptx — no AI image generation.
# Produces typographically consistent, human-quality slides.
# Slide canvas: 13.33" × 7.5" (widescreen 16:9)
# Content area on white slides: x≥0.5", title at y=0.32", content from y≈1.3"

_BODY_CLR  = RGBColor(0x2D, 0x34, 0x36)  # charcoal body text
_MUTED_CLR = RGBColor(0x63, 0x6E, 0x72)  # muted / secondary text
_WHITE_CLR = RGBColor(0xFF, 0xFF, 0xFF)  # pure white
_ROW_ALT   = RGBColor(0xF4, 0xF5, 0xF7)  # light grey for alternating table rows


def _hex_to_rgb(hex_str: str, fallback: Optional[RGBColor] = None) -> RGBColor:
    """Parse '#RRGGBB' or 'RRGGBB' into RGBColor."""
    _fb = fallback if fallback is not None else RGBColor(0x1B, 0x43, 0x32)
    try:
        h = (hex_str or "").lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)) if len(h) == 6 else _fb
    except Exception:
        return _fb


def _tint(color: RGBColor, blend: float = 0.75) -> RGBColor:
    """Blend an RGBColor toward white. blend=0 → same color, blend=1 → pure white."""
    h = str(color)
    r = int(int(h[0:2], 16) + (255 - int(h[0:2], 16)) * blend)
    g = int(int(h[2:4], 16) + (255 - int(h[2:4], 16)) * blend)
    b = int(int(h[4:6], 16) + (255 - int(h[4:6], 16)) * blend)
    return RGBColor(min(255, r), min(255, g), min(255, b))


# ── Clean Deck: individual slide builders ──────────────────────────────────────

def _cdeck_title(prs, primary: RGBColor, title: str, tagline: str,
                 website: str = "", founder: str = "") -> None:
    """Cover slide: primary color bg, white left-aligned title + tagline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, primary)
    _add_textbox(slide, Inches(0.7), Inches(1.9), Inches(11.6), Inches(2.0),
                 title, font_size=52, color=_WHITE_CLR, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri Light")
    if tagline:
        _add_textbox(slide, Inches(0.7), Inches(4.1), Inches(9.5), Inches(0.8),
                     tagline, font_size=22, color=_tint(primary, 0.72),
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")
    footer_parts = [x for x in [website, founder] if x and str(x).strip()]
    if footer_parts:
        _add_textbox(slide, Inches(0.7), Inches(6.7), Inches(11.6), Inches(0.45),
                     "   ·   ".join(footer_parts), font_size=13,
                     color=_tint(primary, 0.60),
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")


def _cdeck_content(prs, primary: RGBColor, title: str, bullets: List[str],
                   subtitle: str = "") -> None:
    """Generic content slide: white bg, title + dash-bulleted list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    if subtitle:
        _add_textbox(slide, Inches(0.5), Inches(1.12), Inches(12.33), Inches(0.45),
                     subtitle, font_size=18, color=_MUTED_CLR,
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")
    content_y = Inches(1.68) if subtitle else Inches(1.32)
    content_h = Inches(5.1) if subtitle else Inches(5.5)
    _add_bullet_textbox(slide, Inches(0.5), content_y, Inches(12.33), content_h,
                        (bullets or [])[:4], font_size=16, color=_BODY_CLR,
                        bullet_char="—", spacing=Pt(16), font_name="Calibri")


def _cdeck_stat_callout(prs, primary: RGBColor, title: str,
                        stats: List[Dict[str, str]]) -> None:
    """Stat callout: 1–3 large numbers centered on white bg."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    stats = (stats or [])[:3]
    n = max(len(stats), 1)
    col_w = Inches(12.33 / n)
    for i, st in enumerate(stats):
        x = Inches(0.5) + col_w * i
        _add_textbox(slide, x, Inches(2.1), col_w, Inches(2.1),
                     st.get("number", ""), font_size=70, color=primary, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")
        _add_textbox(slide, x, Inches(4.25), col_w, Inches(0.75),
                     st.get("label", ""), font_size=16, color=_MUTED_CLR,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")
        if st.get("sublabel"):
            _add_textbox(slide, x, Inches(5.0), col_w, Inches(0.5),
                         st["sublabel"], font_size=13, color=_MUTED_CLR,
                         alignment=PP_ALIGN.CENTER, font_name="Calibri")
        if i < n - 1:
            _add_rect(slide, x + col_w - Inches(0.025), Inches(1.9), Inches(0.025),
                      Inches(3.5), _ROW_ALT)


def _cdeck_two_column(prs, primary: RGBColor, title: str,
                      left_items: List[str], right_items) -> None:
    """Two-column: bullets left | stats (dict list) or bullets (str list) right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    content_y = Inches(1.32)
    _add_bullet_textbox(slide, Inches(0.5), content_y, Inches(5.6), Inches(5.7),
                        [str(x) for x in (left_items or [])][:4],
                        font_size=16, color=_BODY_CLR,
                        bullet_char="—", spacing=Pt(14), font_name="Calibri")
    _add_rect(slide, Inches(6.35), content_y + Inches(0.1), Inches(0.025),
              Inches(5.5), _ROW_ALT)
    right_x = Inches(6.5)
    right_w = Inches(6.33)
    right = right_items or []
    if right and isinstance(right[0], dict):
        for j, item in enumerate(right[:3]):
            iy = content_y + Inches(j * 1.95)
            _add_textbox(slide, right_x, iy, right_w, Inches(1.5),
                         item.get("number", ""), font_size=58, color=primary, bold=True,
                         alignment=PP_ALIGN.LEFT, font_name="Calibri")
            _add_textbox(slide, right_x, iy + Inches(1.15), right_w, Inches(0.55),
                         item.get("label", ""), font_size=14, color=_MUTED_CLR,
                         alignment=PP_ALIGN.LEFT, font_name="Calibri")
    else:
        _add_bullet_textbox(slide, right_x, content_y, right_w, Inches(5.7),
                            [str(x) for x in right][:4],
                            font_size=16, color=_BODY_CLR,
                            bullet_char="—", spacing=Pt(14), font_name="Calibri")


def _cdeck_icon_grid(prs, primary: RGBColor, title: str,
                     items: List[Dict[str, str]]) -> None:
    """Icon grid: 2×2 (4 items) or 3×2 (6 items) — circle + bold label + description."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    items = (items or [])[:6]
    n = len(items)
    cols = 3 if n > 4 else 2
    rows_needed = max((n + cols - 1) // cols, 1)
    grid_y = Inches(1.3)
    cell_w = Inches(12.33 / cols)
    cell_h = Inches(5.9 / rows_needed)
    for i, item in enumerate(items):
        row, col = divmod(i, cols)
        cx = Inches(0.5) + col * cell_w
        cy = grid_y + row * cell_h
        r_size = Inches(0.55)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      cx + (cell_w - r_size) / 2,
                                      cy + Inches(0.2),
                                      r_size, r_size)
        circ.fill.solid()
        circ.fill.fore_color.rgb = primary
        circ.line.fill.background()
        _add_textbox(slide, cx + Inches(0.1), cy + Inches(0.9), cell_w - Inches(0.2),
                     Inches(0.55),
                     item.get("label", ""), font_size=17, color=_BODY_CLR, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")
        _add_textbox(slide, cx + Inches(0.15), cy + Inches(1.5), cell_w - Inches(0.3),
                     Inches(0.9),
                     item.get("description", ""), font_size=13, color=_MUTED_CLR,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")


def _cdeck_flow(prs, primary: RGBColor, title: str,
                steps: List[Dict[str, str]]) -> None:
    """Horizontal numbered flow: 3–5 steps, circles on connector line."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    steps = (steps or [])[:5]
    n = max(len(steps), 1)
    step_w = Inches(12.33 / n)
    line_y = Inches(3.2)
    _add_rect(slide, Inches(0.5) + step_w * 0.35,
              line_y + Inches(0.25),
              Inches(12.33) - step_w * 0.7,
              Inches(0.04), _tint(primary, 0.55))
    for i, step in enumerate(steps):
        sx = Inches(0.5) + i * step_w
        cx = sx + step_w / 2
        c_size = Inches(0.6)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      cx - c_size / 2,
                                      line_y,
                                      c_size, c_size)
        circ.fill.solid()
        circ.fill.fore_color.rgb = primary
        circ.line.fill.background()
        _add_textbox(slide, cx - c_size / 2, line_y + Inches(0.07),
                     c_size, Inches(0.48),
                     str(i + 1), font_size=17, color=_WHITE_CLR, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")
        _add_textbox(slide, sx + Inches(0.1), line_y + Inches(0.75),
                     step_w - Inches(0.2), Inches(0.6),
                     step.get("label", ""), font_size=15, color=_BODY_CLR, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")
        _add_textbox(slide, sx + Inches(0.15), line_y + Inches(1.4),
                     step_w - Inches(0.3), Inches(1.4),
                     step.get("description", ""), font_size=13, color=_MUTED_CLR,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")


def _cdeck_table(prs, primary: RGBColor, title: str,
                 columns: List[str], features: List[Dict[str, Any]]) -> None:
    """Comparison table: primary-color header row, alternating grey data rows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    cols = (columns or [])[:5]
    feats = (features or [])[:8]
    n_cols = len(cols) + 1
    n_rows = len(feats) + 1
    tbl_x = Inches(0.5)
    tbl_y = Inches(1.3)
    tbl_w = Inches(12.33)
    tbl_h = Inches(5.9)
    col_w = tbl_w / n_cols
    row_h = tbl_h / n_rows
    _add_rect(slide, tbl_x, tbl_y, tbl_w, row_h, primary)
    _add_textbox(slide, tbl_x + Inches(0.12), tbl_y, col_w - Inches(0.12), row_h,
                 "Feature", font_size=14, color=_WHITE_CLR, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    for c, cname in enumerate(cols):
        _add_textbox(slide, tbl_x + col_w * (c + 1), tbl_y, col_w, row_h,
                     cname, font_size=14, color=_WHITE_CLR, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")
    for r, feat in enumerate(feats):
        ry = tbl_y + row_h * (r + 1)
        _add_rect(slide, tbl_x, ry, tbl_w, row_h,
                  _ROW_ALT if r % 2 == 0 else _WHITE_CLR)
        _add_textbox(slide, tbl_x + Inches(0.12), ry, col_w - Inches(0.12), row_h,
                     feat.get("feature", ""), font_size=13, color=_BODY_CLR, bold=True,
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")
        for c, val in enumerate((feat.get("values") or [])[:len(cols)]):
            _add_textbox(slide, tbl_x + col_w * (c + 1), ry, col_w, row_h,
                         str(val), font_size=13, color=_BODY_CLR,
                         alignment=PP_ALIGN.CENTER, font_name="Calibri")


def _cdeck_timeline(prs, primary: RGBColor, title: str,
                    milestones: List[Dict[str, str]]) -> None:
    """Timeline: horizontal line with alternating above/below milestone labels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, _WHITE_CLR)
    _add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12.33), Inches(0.78),
                 title, font_size=36, color=primary, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    mss = (milestones or [])[:6]
    n = max(len(mss), 1)
    tl_x = Inches(0.7)
    tl_y = Inches(3.9)
    tl_w = Inches(11.6)
    _add_rect(slide, tl_x, tl_y, tl_w, Inches(0.05), primary)
    step = tl_w / (n - 1) if n > 1 else tl_w
    for i, ms in enumerate(mss):
        mx = tl_x + step * i if n > 1 else tl_x + tl_w / 2
        item_w = Inches(min(1.9, 11.6 / n * 0.85))
        d = Inches(0.22)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     mx - d / 2,
                                     tl_y - d / 2 + Inches(0.025),
                                     d, d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = primary
        dot.line.fill.background()
        hw = item_w / 2
        if i % 2 == 0:
            _add_textbox(slide, mx - hw, tl_y - Inches(1.7), item_w, Inches(0.45),
                         ms.get("date", ""), font_size=12, color=_MUTED_CLR,
                         alignment=PP_ALIGN.CENTER, font_name="Calibri")
            _add_textbox(slide, mx - hw, tl_y - Inches(1.2), item_w, Inches(0.55),
                         ms.get("label", ""), font_size=14, color=_BODY_CLR, bold=True,
                         alignment=PP_ALIGN.CENTER, font_name="Calibri")
            if ms.get("description"):
                _add_textbox(slide, mx - hw, tl_y - Inches(0.62), item_w, Inches(0.5),
                             ms["description"], font_size=12, color=_MUTED_CLR,
                             alignment=PP_ALIGN.CENTER, font_name="Calibri")
        else:
            _add_textbox(slide, mx - hw, tl_y + Inches(0.25), item_w, Inches(0.55),
                         ms.get("label", ""), font_size=14, color=_BODY_CLR, bold=True,
                         alignment=PP_ALIGN.CENTER, font_name="Calibri")
            _add_textbox(slide, mx - hw, tl_y + Inches(0.8), item_w, Inches(0.45),
                         ms.get("date", ""), font_size=12, color=_MUTED_CLR,
                         alignment=PP_ALIGN.CENTER, font_name="Calibri")
            if ms.get("description"):
                _add_textbox(slide, mx - hw, tl_y + Inches(1.25), item_w, Inches(0.5),
                             ms["description"], font_size=12, color=_MUTED_CLR,
                             alignment=PP_ALIGN.CENTER, font_name="Calibri")


def _cdeck_closing(prs, primary: RGBColor, title: str,
                   tagline: str = "", contact: str = "", cta: str = "") -> None:
    """Closing slide: mirrors title — primary bg, white left-aligned text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, primary)
    _add_textbox(slide, Inches(0.7), Inches(1.8), Inches(11.6), Inches(1.9),
                 title, font_size=46, color=_WHITE_CLR, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri Light")
    if tagline:
        _add_textbox(slide, Inches(0.7), Inches(3.85), Inches(9.5), Inches(0.75),
                     tagline, font_size=22, color=_tint(primary, 0.72),
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")
    if cta:
        _add_textbox(slide, Inches(0.7), Inches(4.85), Inches(7.0), Inches(0.7),
                     cta, font_size=22, color=_WHITE_CLR, bold=True,
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")
    if contact:
        _add_textbox(slide, Inches(0.7), Inches(6.65), Inches(11.6), Inches(0.5),
                     contact, font_size=14, color=_tint(primary, 0.60),
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")


# ── Body Converters (fallback parsers when structured data not provided) ───────

def _body_to_stats(body: List[str]) -> List[Dict[str, str]]:
    out = []
    for line in (body or [])[:3]:
        parts = (line or "").strip().split(" ", 1)
        out.append({"number": parts[0], "label": parts[1] if len(parts) > 1 else ""})
    return out or [{"number": "—", "label": ""}]


def _body_to_items(body: List[str]) -> List[Dict[str, str]]:
    out = []
    for line in (body or [])[:6]:
        parts = (line or "").split(":", 1)
        out.append({"label": parts[0].strip(),
                    "description": parts[1].strip() if len(parts) > 1 else ""})
    return out


def _body_to_steps(body: List[str]) -> List[Dict[str, str]]:
    out = []
    for line in (body or [])[:5]:
        parts = (line or "").split(":", 1)
        out.append({"label": parts[0].strip(),
                    "description": parts[1].strip() if len(parts) > 1 else ""})
    return out


def _body_to_features(body: List[str]) -> List[Dict[str, Any]]:
    return [{"feature": (line or "").strip(), "values": ["✓", "✓", "✓"]}
            for line in (body or [])[:8]]


def _body_to_milestones(body: List[str]) -> List[Dict[str, str]]:
    out = []
    for i, line in enumerate((body or [])[:6]):
        parts = (line or "").split(":", 1)
        out.append({"label": parts[0].strip(),
                    "description": parts[1].strip() if len(parts) > 1 else "",
                    "date": f"Phase {i + 1}"})
    return out


# ── Main Clean Deck Builder ─────────────────────────────────────────────────────

async def create_clean_deck_async(
    topic: str,
    slides_plan: List[Dict[str, Any]],
    business_name: str = "My Business",
    brand_color: str = "",
) -> Dict[str, Any]:
    """
    Build a clean python-pptx presentation using the structured design system.
    No AI image generation — entirely rendered via python-pptx for crisp,
    typographically consistent, human-quality output.

    Each slide in slides_plan may include:
        title        (str)
        layout       (str): title | stat_callout | two_column | icon_grid | flow |
                            comparison_table | timeline | closing | content
        body         (list[str])  : bullets (fallback if structured data absent)
        subtitle     (str)        : subheader line on content slides
        tagline      (str)        : tagline for title / closing slides
        stats        (list[dict]) : [{number, label, sublabel}] for stat_callout
        items        (list[dict]) : [{label, description}] for icon_grid
        steps        (list[dict]) : [{label, description}] for flow
        columns      (list[str])  : tier names for comparison_table
        features     (list[dict]) : [{feature, values[]}] for comparison_table
        milestones   (list[dict]) : [{date, label, description}] for timeline
        left_items   (list[str])  : left column bullets for two_column
        right_items  (list)       : right column — str list OR [{number, label}]
        website      (str)        : for title/closing footer
        founder      (str)        : for title slide footer
        contact      (str)        : for closing footer
        cta          (str)        : call-to-action text for closing slide
    """
    import logging as _log

    primary = _hex_to_rgb(brand_color) if brand_color else RGBColor(0x1B, 0x43, 0x32)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, sd in enumerate(slides_plan):
        layout = ((sd.get("layout") or "") or ("title" if i == 0 else "content")).lower().strip()
        body   = sd.get("body") or []

        if layout == "title":
            _cdeck_title(prs, primary,
                         sd.get("title", topic),
                         sd.get("tagline") or sd.get("subtitle") or (body[0] if body else ""),
                         sd.get("website", ""),
                         sd.get("founder", ""))

        elif layout == "stat_callout":
            _cdeck_stat_callout(prs, primary, sd.get("title", ""),
                                sd.get("stats") or _body_to_stats(body))

        elif layout == "two_column":
            left  = sd.get("left_items") or sd.get("left") or []
            right = sd.get("right_items") or sd.get("right") or []
            if not left and not right and body:
                half  = max(1, len(body) // 2)
                left, right = body[:half], body[half:]
            elif not left:
                left = body
            _cdeck_two_column(prs, primary, sd.get("title", ""), left, right)

        elif layout == "icon_grid":
            _cdeck_icon_grid(prs, primary, sd.get("title", ""),
                             sd.get("items") or _body_to_items(body))

        elif layout == "flow":
            _cdeck_flow(prs, primary, sd.get("title", ""),
                        sd.get("steps") or _body_to_steps(body))

        elif layout == "comparison_table":
            _cdeck_table(prs, primary, sd.get("title", ""),
                         sd.get("columns") or ["Starter", "Pro", "Enterprise"],
                         sd.get("features") or _body_to_features(body))

        elif layout == "timeline":
            _cdeck_timeline(prs, primary, sd.get("title", ""),
                            sd.get("milestones") or _body_to_milestones(body))

        elif layout == "closing":
            _cdeck_closing(prs, primary,
                           sd.get("title", ""),
                           sd.get("tagline") or sd.get("subtitle") or "",
                           sd.get("contact", ""),
                           sd.get("cta") or "Let's talk.")

        else:  # content fallback
            _cdeck_content(prs, primary, sd.get("title", ""),
                           body, sd.get("subtitle", ""))

    filename = f"clean-deck-{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(PRESENTATIONS_DIR, filename)
    prs.save(filepath)

    file_url = f"/api/media/presentations/{filename}"
    try:
        from pathlib import Path as _Path
        from image_handler import S3Handler
        _fbytes = _Path(filepath).read_bytes()
        _b64    = base64.b64encode(_fbytes).decode()
        s3_name = f"clean-deck-{uuid.uuid4().hex[:8]}.pptx"
        s3_url  = await S3Handler.upload_file(
            _b64, s3_name,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        if s3_url:
            file_url = s3_url
    except Exception as e:
        _log.warning("[clean_deck] S3 upload skipped: %s", e)
    finally:
        try:
            os.unlink(filepath)
        except Exception:
            pass

    return {
        "success":     True,
        "url":         file_url,
        "filename":    filename,
        "slide_count": len(slides_plan),
        "deck_type":   "clean",
    }
